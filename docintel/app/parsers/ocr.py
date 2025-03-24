import os
import base64
import tempfile
import re
import ssl
import certifi
import time
from typing import List, Dict, Any, Optional, BinaryIO, Tuple
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from openai import AzureOpenAI
from app.utils.logging import log_step, Timer
from concurrent.futures import ThreadPoolExecutor, as_completed
import httpx


class OCRProcessor:
    """
    Processes complex documents using OCR with Azure OpenAI's GPT-4o-mini.
    """
    
    def __init__(self, max_workers: int = 20):
        """
        Initialize the OCR processor.
        
        Args:
            max_workers: Maximum number of parallel workers for OCR processing
        """
        self.max_workers = max_workers
        
        # Configure SSL context for Azure OpenAI
        self.ssl_context = ssl.create_default_context(cafile=certifi.where())
        self.ssl_context.verify_mode = ssl.CERT_REQUIRED
        self.ssl_context.check_hostname = True
    
    def process_file(
        self, 
        file_path: str,
        file_type: str
    ) -> List[Dict[str, Any]]:
        """
        Process a file using OCR.
        
        Args:
            file_path: Path to the file
            file_type: Type of file (pdf, docx, pptx)
            
        Returns:
            List of dictionaries with extracted text and metadata for each page
        """
        with Timer("OCR Processing"):
            log_step("OCR", f"Processing file: {file_path}")
            
            # Extract images based on file type
            images = self._extract_images(file_path, file_type)
            log_step("OCR", f"Extracted {len(images)} images from {file_type.upper()}")
            
            # Process images in parallel
            results = []
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                # Create a dictionary to map futures to their page numbers
                future_to_page = {
                    executor.submit(self._process_single_image, image_data, i + 1): i + 1
                    for i, image_data in enumerate(images)
                }
                
                # Process completed futures as they finish
                for future in as_completed(future_to_page):
                    page_num = future_to_page[future]
                    try:
                        text = future.result()
                        if text:  # Only add non-empty results
                            results.append({
                                "page_number": page_num,
                                "text": text,
                                "is_ocr": True
                            })
                        else:
                            log_step("OCR", f"No text extracted from page/slide {page_num}", level="warning")
                    except Exception as e:
                        log_step("OCR", f"Error processing page/slide {page_num}: {str(e)}", level="error")
            
            # Sort results by page number
            results.sort(key=lambda x: x["page_number"])
            log_step("OCR", f"Completed OCR processing with {len(results)} pages/slides extracted")
            return results
    
    def process_stream(
        self, 
        file_stream: BinaryIO,
        file_type: str,
        filename: str
    ) -> List[Dict[str, Any]]:
        """
        Process a file stream using OCR.
        
        Args:
            file_stream: File-like object
            file_type: Type of file (pdf, docx, pptx)
            filename: Original filename
            
        Returns:
            List of dictionaries with extracted text and metadata for each page
        """
        with Timer("OCR Processing"):
            log_step("OCR", f"Processing file stream: {filename}")
            
            # Save stream to temporary file
            temp_file_path = f"temp_{os.path.basename(filename)}"
            with open(temp_file_path, "wb") as f:
                f.write(file_stream.read())
                
            try:
                # Extract images based on file type
                images = self._extract_images(temp_file_path, file_type)
                log_step("OCR", f"Extracted {len(images)} images from {file_type.upper()}")
                
                # Process images in parallel
                results = []
                with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                    # Create a dictionary to map futures to their page numbers
                    future_to_page = {
                        executor.submit(self._process_single_image, image_data, i + 1): i + 1
                        for i, image_data in enumerate(images)
                    }
                    
                    # Process completed futures as they finish
                    for future in as_completed(future_to_page):
                        page_num = future_to_page[future]
                        try:
                            text = future.result()
                            if text:  # Only add non-empty results
                                results.append({
                                    "page_number": page_num,
                                    "text": text,
                                    "is_ocr": True
                                })
                            else:
                                log_step("OCR", f"No text extracted from page/slide {page_num}", level="warning")
                        except Exception as e:
                            log_step("OCR", f"Error processing page/slide {page_num}: {str(e)}", level="error")
                
                # Sort results by page number
                results.sort(key=lambda x: x["page_number"])
                log_step("OCR", f"Completed OCR processing with {len(results)} pages/slides extracted")
                return results
            finally:
                # Clean up temporary file
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
    
    def _process_single_image(self, image_data: bytes, page_num: int) -> Optional[str]:
        """
        Process a single image with OCR.
        
        Args:
            image_data: Image data as bytes
            page_num: Page number for logging
            
        Returns:
            Extracted text or None if extraction failed
        """
        if not image_data:
            log_step("OCR", f"No image data for page/slide {page_num}", level="warning")
            return None
            
        log_step("OCR", f"Processing page/slide {page_num}")
        try:
            return self._perform_llm_ocr(image_data)
        except Exception as e:
            log_step("OCR", f"Error in OCR for page/slide {page_num}: {str(e)}", level="error")
            return None
    
    def _extract_images(self, file_path: str, file_type: str) -> List[bytes]:
        """
        Extract images from a file based on its type.
        
        Args:
            file_path: Path to the file
            file_type: Type of file (pdf, docx, pptx)
            
        Returns:
            List of image data as bytes
        """
        images = []
        
        if file_type == "pdf":
            images = self._extract_images_from_pdf(file_path)
        elif file_type == "docx":
            images = self._extract_images_from_docx(file_path)
        elif file_type == "pptx":
            images = self._extract_images_from_pptx(file_path)
        else:
            log_step("OCR", f"Unsupported file type for OCR: {file_type}", level="error")
            
        return images
    
    def _extract_images_from_pdf(self, file_path: str) -> List[bytes]:
        """
        Extract images from a PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of image data as bytes
        """
        images = []
        
        try:
            # Open PDF
            pdf = fitz.open(file_path)
            
            # Process each page
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                
                # Render page as image with higher resolution for better OCR
                # Use a zoom factor of 2 for better quality
                zoom = 2.0
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Convert to PIL Image for potential processing
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # Save to bytes
                buffer = BytesIO()
                img.save(buffer, format="JPEG", quality=95)
                images.append(buffer.getvalue())
                
            log_step("OCR", f"Extracted {len(images)} images from PDF")
        except Exception as e:
            log_step("OCR", f"Error extracting images from PDF: {str(e)}", level="error")
        
        return images
    
    def _extract_images_from_docx(self, file_path: str) -> List[bytes]:
        """
        Extract images from a DOCX file by rendering each page.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            List of image data as bytes
        """
        images = []
        
        try:
            # Open the DOCX file
            doc = docx.Document(file_path)
            
            # We'll create a visual representation of each page
            # Since docx doesn't have a direct "page" concept, we'll estimate
            
            # First attempt: convert to PDF using a third-party tool if available
            pdf_path = None
            try:
                # Try to use docx2pdf if installed
                from docx2pdf import convert
                
                # Create temporary PDF file
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
                    pdf_path = temp_pdf.name
                
                # Convert DOCX to PDF
                convert(file_path, pdf_path)
                
                # If conversion successful, extract images from the PDF
                if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                    return self._extract_images_from_pdf(pdf_path)
            except ImportError:
                log_step("OCR", "docx2pdf not available, using fallback method", level="warning")
            finally:
                # Clean up temporary PDF file
                if pdf_path and os.path.exists(pdf_path):
                    os.remove(pdf_path)
            
            # Fallback: Create image for each paragraph/section
            # This is a simplistic approach - production would use better rendering
            
            # Create a single image with text content
            # Width and height for the image
            width, height = 1000, 1400
            margin = 50
            line_height = 25
            
            # Create a new image for each approximate "page" of content
            current_y = margin
            page_num = 1
            page_image = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(page_image)
            
            # Try to load a font, fall back to default if not available
            try:
                font = ImageFont.truetype("Arial", 16)
            except IOError:
                font = ImageFont.load_default()
            
            # Extract text from paragraphs and tables
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    para = docx.Document(file_path)._element[0].xpath('//w:p')[0]
                    text = "".join(t.text for t in para.xpath('.//w:t'))
                    
                    # Wrap text to fit width
                    words = text.split()
                    line = ""
                    
                    for word in words:
                        test_line = line + word + " "
                        line_width = draw.textlength(test_line, font=font)
                        
                        if line_width > width - 2 * margin:
                            draw.text((margin, current_y), line, fill="black", font=font)
                            current_y += line_height
                            line = word + " "
                        else:
                            line = test_line
                    
                    # Draw remaining text
                    if line:
                        draw.text((margin, current_y), line, fill="black", font=font)
                        current_y += line_height * 1.5  # Extra space after paragraph
                
                elif element.tag.endswith('tbl'):  # Table
                    # Simplified table representation
                    draw.text((margin, current_y), "[TABLE CONTENT]", fill="black", font=font)
                    current_y += line_height * 2
                
                # Check if we need a new page
                if current_y > height - margin:
                    # Save current page
                    buffer = BytesIO()
                    page_image.save(buffer, format="JPEG", quality=95)
                    images.append(buffer.getvalue())
                    
                    # Create new page
                    page_num += 1
                    page_image = Image.new('RGB', (width, height), color='white')
                    draw = ImageDraw.Draw(page_image)
                    current_y = margin
            
            # Save the last page
            if current_y > margin:  # Only if there's content
                buffer = BytesIO()
                page_image.save(buffer, format="JPEG", quality=95)
                images.append(buffer.getvalue())
            
            log_step("OCR", f"Extracted {len(images)} images from DOCX")
        except Exception as e:
            log_step("OCR", f"Error extracting images from DOCX: {str(e)}", level="error")
            
            # Fallback to a simple representation
            img = Image.new('RGB', (800, 1000), color='white')
            draw = ImageDraw.Draw(img)
            
            try:
                # Try to extract some text for the placeholder
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs[:20]])
                
                # Add text to image
                lines = text.split('\n')
                y_position = 50
                for line in lines:
                    draw.text((50, y_position), line, fill="black")
                    y_position += 30
                
                draw.text((50, y_position + 30), "... [Document continues] ...", fill="black")
            except:
                # If all else fails, create a blank image with message
                draw.text((50, 50), "DOCX CONTENT", fill="black")
                draw.text((50, 100), "The content couldn't be rendered as an image.", fill="black")
                draw.text((50, 150), "Using OCR service to extract text directly.", fill="black")
            
            buffer = BytesIO()
            img.save(buffer, format="JPEG", quality=95)
            images.append(buffer.getvalue())
        
        return images
    
    def _extract_images_from_pptx(self, file_path: str) -> List[bytes]:
        """
        Extract images from a PPTX file by rendering each slide.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of image data as bytes
        """
        images = []
        
        try:
            # Open the PPTX file
            presentation = Presentation(file_path)
            
            # Try to use python-pptx-export if available (better rendering)
            try:
                from pptx_export import export_slides
                
                # Create temporary directory for slides
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Export slides as images
                    export_slides(file_path, temp_dir, "jpg")
                    
                    # Read exported images in order
                    slide_files = sorted([f for f in os.listdir(temp_dir) if f.endswith('.jpg')], 
                                        key=lambda x: int(re.search(r'slide(\d+)', x).group(1)))
                    
                    for slide_file in slide_files:
                        with open(os.path.join(temp_dir, slide_file), 'rb') as f:
                            images.append(f.read())
            
            except ImportError:
                log_step("OCR", "pptx_export not available, using fallback method", level="warning")
                
                # Fallback: Create images with slide content
                for i, slide in enumerate(presentation.slides):
                    # Create image for slide
                    width, height = 1280, 720  # 16:9 aspect ratio
                    img = Image.new('RGB', (width, height), color='white')
                    draw = ImageDraw.Draw(img)
                    
                    # Try to load a font, fall back to default if not available
                    try:
                        title_font = ImageFont.truetype("Arial", 36)
                        content_font = ImageFont.truetype("Arial", 24)
                    except IOError:
                        title_font = ImageFont.load_default()
                        content_font = ImageFont.load_default()
                    
                    # Add slide number
                    draw.text((width - 100, height - 50), f"Slide {i+1}", fill="black", font=content_font)
                    
                    # Extract and add slide title
                    title_shape = None
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.name.startswith('Title'):
                            title_shape = shape
                            break
                    
                    if title_shape and title_shape.has_text_frame:
                        title_text = title_shape.text
                        draw.text((50, 50), title_text, fill="black", font=title_font)
                    
                    # Extract and add content
                    y_position = 150
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != title_shape:
                            text = shape.text
                            # Simple text wrapping
                            words = text.split()
                            line = ""
                            for word in words:
                                test_line = line + word + " "
                                line_width = draw.textlength(test_line, font=content_font)
                                
                                if line_width > width - 100:
                                    draw.text((50, y_position), line, fill="black", font=content_font)
                                    y_position += 30
                                    line = word + " "
                                else:
                                    line = test_line
                            
                            if line:
                                draw.text((50, y_position), line, fill="black", font=content_font)
                            y_position += 50
                    
                    # Save the slide image
                    buffer = BytesIO()
                    img.save(buffer, format="JPEG", quality=95)
                    images.append(buffer.getvalue())
            
            log_step("OCR", f"Extracted {len(images)} images from PPTX")
        except Exception as e:
            log_step("OCR", f"Error extracting images from PPTX: {str(e)}", level="error")
            
            # Fallback to a simple representation
            img = Image.new('RGB', (1280, 720), color='white')
            draw = ImageDraw.Draw(img)
            
            try:
                # Try to extract slide titles for the placeholder
                presentation = Presentation(file_path)
                titles = []
                for i, slide in enumerate(presentation.slides[:10]):
                    title = "Untitled Slide"
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.name.startswith('Title'):
                            title = shape.text
                            break
                    titles.append(f"Slide {i+1}: {title}")
                
                # Add text to image
                y_position = 50
                for title in titles:
                    draw.text((50, y_position), title, fill="black")
                    y_position += 30
                
                if len(presentation.slides) > 10:
                    draw.text((50, y_position + 30), f"... [{len(presentation.slides) - 10} more slides] ...", fill="black")
            except:
                # If all else fails, create a blank image with message
                draw.text((50, 50), "PPTX CONTENT", fill="black")
                draw.text((50, 100), "The slides couldn't be rendered as images.", fill="black")
                draw.text((50, 150), "Using OCR service to extract text directly.", fill="black")
            
            buffer = BytesIO()
            img.save(buffer, format="JPEG", quality=95)
            images.append(buffer.getvalue())
        
        return images
    
    def _perform_llm_ocr(self, image_data: bytes) -> Optional[str]:
        """
        Perform OCR using Azure OpenAI's vision model.
        
        Args:
            image_data: Image data as bytes
            
        Returns:
            Extracted text or None if extraction failed
        """
        if not image_data:
            log_step("LLM OCR", "No image data provided", level="warning")
            return None
            
        try:
            # Convert image to base64
            img_base64 = base64.b64encode(image_data).decode('utf-8')
            
            # Get Azure OpenAI configuration
            api_key = os.getenv("AZURE_OPENAI_API_KEY")
            api_base = os.getenv("AZURE_OPENAI_API_BASE")
            deployment = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME")
            api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
            
            if not all([api_key, api_base, deployment]):
                log_step("LLM OCR", "Missing Azure OpenAI configuration", level="error")
                return None
            
            # Initialize Azure OpenAI client with SSL context and custom http_client to avoid proxies issue
            http_client = httpx.Client(verify=self.ssl_context)
            client = AzureOpenAI(
                api_key=api_key,
                api_version=api_version,
                azure_endpoint=api_base,
                http_client=http_client
            )

            # Create message format according to the vision API requirements
            messages = [
                {
                    "role": "system",
                    "content": [
                        {
                            "type": "text",
                            "text": "You are an OCR assistant that extracts text from document images. Extract ALL text, preserve formatting, and use Markdown when appropriate."
                        }
                    ]
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract all the text from this document image. Preserve the layout with Markdown formatting. If you see tables, format them as Markdown tables."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{img_base64}"
                            }
                        }
                    ]
                }
            ]

            # Call Azure OpenAI for OCR with retry logic
            max_retries = 3
            retry_delay = 1  # seconds
            
            for attempt in range(max_retries):
                try:
                    log_step("LLM OCR", f"Sending image to Azure OpenAI model: {deployment} (attempt {attempt + 1})")
                    completion = client.chat.completions.create(
                        model=deployment,
                        messages=messages,
                        max_tokens=4096,
                        temperature=0.0,  # Use 0 for most accurate OCR
                        stream=False
                    )
                    
                    # Extract text from response
                    if completion and completion.choices:
                        extracted_text = completion.choices[0].message.content
                        if extracted_text:
                            log_step("LLM OCR", f"Extracted {len(extracted_text)} characters from image")
                            return extracted_text
                    
                    log_step("LLM OCR", "No text extracted from response", level="warning")
                    return None
                    
                except Exception as e:
                    if attempt < max_retries - 1:
                        log_step("LLM OCR", f"Attempt {attempt + 1} failed: {str(e)}, retrying...", level="warning")
                        time.sleep(retry_delay * (attempt + 1))  # Exponential backoff
                    else:
                        log_step("LLM OCR", f"All attempts failed: {str(e)}", level="error")
                        return None
                        
        except Exception as e:
            log_step("LLM OCR", f"Error: {str(e)}", level="error")
            return None