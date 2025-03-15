import os
import time
from typing import Dict, List, Any, Optional, BinaryIO, Tuple
import pptx
from io import BytesIO

from app.parsers.base_parser import BaseDocumentParser
from app.parsers.ocr import OCRProcessor
from app.chunking.models import ProcessedDocument
from app.chunking.chunker import DocumentChunker
from app.utils.logging import log_step, Timer


class PPTXParser(BaseDocumentParser):
    """Parser for PPTX documents."""
    
    def __init__(self, chunker: Optional[DocumentChunker] = None):
        """
        Initialize PPTX parser.
        
        Args:
            chunker: Document chunker instance (optional)
        """
        super().__init__()
        self.chunker = chunker or DocumentChunker()
        self.ocr_processor = OCRProcessor()
    
    def parse(
        self, 
        file_path: str, 
        filename: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """
        Parse a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            filename: Original filename (uses basename of file_path if not provided)
            metadata: Additional metadata
            
        Returns:
            ProcessedDocument object with extracted content and metadata
        """
        with Timer("PPTX Parsing"):
            # Get filename if not provided
            if not filename:
                filename = os.path.basename(file_path)
                
            log_step("PPTX Parsing", f"Parsing PPTX file: {filename}")
            
            # Get file size
            file_size = os.path.getsize(file_path)
            
            # Prepare metadata
            doc_metadata = self.prepare_metadata(filename, file_size, metadata)
            
            # Start timer for processing
            start_time = time.time()
            
            try:
                # Try simple PPTX parsing first
                text_by_slide, is_complex = self._extract_text_from_pptx(file_path)
                
                # If simple parsing failed or detected a complex document, use OCR
                if is_complex:
                    log_step("PPTX Parsing", "Complex PPTX detected, using OCR")
                    ocr_results = self.ocr_processor.process_file(file_path, "pptx")
                    
                    # Convert OCR results to text by slide format
                    text_by_slide = {}
                    if ocr_results:
                        text_by_slide = {
                            result["page_number"]: {
                                "text": result["text"],
                                "is_ocr": True
                            }
                            for result in ocr_results
                            if result and result.get("text")
                        }
                    
                    # If OCR failed to extract any text, try one more time with higher quality
                    if not text_by_slide:
                        log_step("PPTX Parsing", "OCR failed, retrying with higher quality settings", level="warning")
                        self.ocr_processor = OCRProcessor(max_workers=2)  # Reduce workers but increase quality
                        ocr_results = self.ocr_processor.process_file(file_path, "pptx")
                        if ocr_results:
                            text_by_slide = {
                                result["page_number"]: {
                                    "text": result["text"],
                                    "is_ocr": True
                                }
                                for result in ocr_results
                                if result and result.get("text")
                            }
                
                # No text extracted
                if not text_by_slide:
                    log_step("PPTX Parsing", "No text extracted from PPTX", level="warning")
                    return ProcessedDocument(
                        document_id=self.document_id,
                        filename=filename,
                        file_type="pptx",
                        file_size=file_size,
                        total_pages=0,
                        total_chunks=0,
                        chunks=[],
                        processing_time=time.time() - start_time,
                        is_complex=is_complex
                    )
                
                # Process each slide
                all_chunks = []
                
                for slide_num, slide_data in text_by_slide.items():
                    slide_text = slide_data["text"]
                    is_ocr = slide_data.get("is_ocr", False)
                    
                    # Skip empty slides
                    if not slide_text:
                        continue
                    
                    # Add slide number to metadata
                    slide_metadata = doc_metadata.copy()
                    slide_metadata["slide_number"] = slide_num
                    slide_metadata["is_ocr"] = is_ocr
                    
                    # Chunk the slide
                    slide_chunks = self.chunker.chunk_document(
                        slide_text,
                        slide_metadata,
                        use_headings=True,
                        is_ocr=is_ocr
                    )
                    
                    all_chunks.extend(slide_chunks)
                
                # Calculate total number of slides
                total_slides = max(text_by_slide.keys()) if text_by_slide else 0
                
                # Create processed document
                processed_doc = ProcessedDocument(
                    document_id=self.document_id,
                    filename=filename,
                    file_type="pptx",
                    file_size=file_size,
                    total_pages=total_slides,
                    total_chunks=len(all_chunks),
                    chunks=all_chunks,
                    processing_time=time.time() - start_time,
                    is_complex=is_complex
                )
                
                log_step("PPTX Parsing", f"Completed parsing PPTX with {total_slides} slides and {len(all_chunks)} chunks")
                return processed_doc
                
            except Exception as e:
                log_step("PPTX Parsing", f"Error parsing PPTX: {str(e)}", level="error")
                raise
    
    def parse_stream(
        self, 
        file_stream: BinaryIO,
        filename: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """
        Parse a PPTX from a file stream.
        
        Args:
            file_stream: File-like object containing the PPTX
            filename: Original filename
            metadata: Additional metadata
            
        Returns:
            ProcessedDocument object with extracted content and metadata
        """
        with Timer("PPTX Stream Parsing"):
            log_step("PPTX Parsing", f"Parsing PPTX from stream: {filename}")
            
            # Get file size
            file_stream.seek(0, os.SEEK_END)
            file_size = file_stream.tell()
            file_stream.seek(0)
            
            # Prepare metadata
            doc_metadata = self.prepare_metadata(filename, file_size, metadata)
            
            # Start timer for processing
            start_time = time.time()
            
            try:
                # Create in-memory file
                pptx_data = file_stream.read()
                memory_stream = BytesIO(pptx_data)
                
                # Try simple PPTX parsing first
                text_by_slide, is_complex = self._extract_text_from_pptx_stream(memory_stream)
                
                # If simple parsing failed or detected a complex document, use OCR
                if is_complex:
                    log_step("PPTX Parsing", "Complex PPTX detected, using OCR")
                    # Reset file stream
                    file_stream.seek(0)
                    ocr_results = self.ocr_processor.process_stream(file_stream, "pptx", filename)
                    
                    # Convert OCR results to text by slide format
                    text_by_slide = {}
                    if ocr_results:
                        text_by_slide = {
                            result["page_number"]: {
                                "text": result["text"],
                                "is_ocr": True
                            }
                            for result in ocr_results
                            if result and result.get("text")
                        }
                    
                    # If OCR failed to extract any text, try one more time with higher quality
                    if not text_by_slide:
                        log_step("PPTX Parsing", "OCR failed, retrying with higher quality settings", level="warning")
                        self.ocr_processor = OCRProcessor(max_workers=2)  # Reduce workers but increase quality
                        ocr_results = self.ocr_processor.process_stream(file_stream, "pptx", filename)
                        if ocr_results:
                            text_by_slide = {
                                result["page_number"]: {
                                    "text": result["text"],
                                    "is_ocr": True
                                }
                                for result in ocr_results
                                if result and result.get("text")
                            }
                
                # No text extracted
                if not text_by_slide:
                    log_step("PPTX Parsing", "No text extracted from PPTX", level="warning")
                    return ProcessedDocument(
                        document_id=self.document_id,
                        filename=filename,
                        file_type="pptx",
                        file_size=file_size,
                        total_pages=0,
                        total_chunks=0,
                        chunks=[],
                        processing_time=time.time() - start_time,
                        is_complex=is_complex
                    )
                
                # Process each slide
                all_chunks = []
                
                for slide_num, slide_data in text_by_slide.items():
                    slide_text = slide_data["text"]
                    is_ocr = slide_data.get("is_ocr", False)
                    
                    # Skip empty slides
                    if not slide_text:
                        continue
                    
                    # Add slide number to metadata
                    slide_metadata = doc_metadata.copy()
                    slide_metadata["slide_number"] = slide_num
                    slide_metadata["is_ocr"] = is_ocr
                    
                    # Chunk the slide
                    slide_chunks = self.chunker.chunk_document(
                        slide_text,
                        slide_metadata,
                        use_headings=True,
                        is_ocr=is_ocr
                    )
                    
                    all_chunks.extend(slide_chunks)
                
                # Calculate total number of slides
                total_slides = max(text_by_slide.keys()) if text_by_slide else 0
                
                # Create processed document
                processed_doc = ProcessedDocument(
                    document_id=self.document_id,
                    filename=filename,
                    file_type="pptx",
                    file_size=file_size,
                    total_pages=total_slides,
                    total_chunks=len(all_chunks),
                    chunks=all_chunks,
                    processing_time=time.time() - start_time,
                    is_complex=is_complex
                )
                
                log_step("PPTX Parsing", f"Completed parsing PPTX from stream with {total_slides} slides and {len(all_chunks)} chunks")
                return processed_doc
                
            except Exception as e:
                log_step("PPTX Parsing", f"Error parsing PPTX stream: {str(e)}", level="error")
                raise
    
    def _extract_text_from_pptx(self, file_path: str) -> Tuple[Dict[int, Dict[str, Any]], bool]:
        """
        Extract text from PPTX using python-pptx.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Tuple of (text by slide dictionary, is complex document flag)
        """
        text_by_slide = {}
        is_complex = False
        
        try:
            # Load presentation
            presentation = pptx.Presentation(file_path)
            
            # Extract text from each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract title if available
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    if title_text:
                        slide_text.append(f"# {title_text}")
                
                # Extract text from each shape
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        slide_text.append(shape.text)
                
                # Join all text from slide
                full_text = "\n\n".join(slide_text)
                
                # Check if slide has meaningful text
                if not full_text or len(full_text) < 20:
                    # This might be a slide with mainly images or charts
                    # Flag as potentially complex
                    is_complex = True
                
                text_by_slide[slide_num] = {
                    "text": full_text,
                    "is_ocr": False
                }
            
            # If no text found in presentation, flag as complex
            total_text = sum(len(slide_data["text"]) for slide_data in text_by_slide.values())
            if total_text < 100:
                is_complex = True
            
            return text_by_slide, is_complex
            
        except Exception as e:
            log_step("PPTX Parsing", f"Error in simple PPTX parsing: {str(e)}", level="warning")
            return {}, True
    
    def _extract_text_from_pptx_stream(self, memory_stream: BytesIO) -> Tuple[Dict[int, Dict[str, Any]], bool]:
        """
        Extract text from PPTX stream using python-pptx.
        
        Args:
            memory_stream: BytesIO object containing the PPTX
            
        Returns:
            Tuple of (text by slide dictionary, is complex document flag)
        """
        text_by_slide = {}
        is_complex = False
        
        try:
            # Reset stream position
            memory_stream.seek(0)
            
            # Load presentation
            presentation = pptx.Presentation(memory_stream)
            
            # Extract text from each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract title if available
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    if title_text:
                        slide_text.append(f"# {title_text}")
                
                # Extract text from each shape
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        slide_text.append(shape.text)
                
                # Join all text from slide
                full_text = "\n\n".join(slide_text)
                
                # Check if slide has meaningful text
                if not full_text or len(full_text) < 20:
                    # This might be a slide with mainly images or charts
                    # Flag as potentially complex
                    is_complex = True
                
                text_by_slide[slide_num] = {
                    "text": full_text,
                    "is_ocr": False
                }
            
            # If no text found in presentation, flag as complex
            total_text = sum(len(slide_data["text"]) for slide_data in text_by_slide.values())
            if total_text < 100:
                is_complex = True
            
            return text_by_slide, is_complex
            
        except Exception as e:
            log_step("PPTX Parsing", f"Error in simple PPTX parsing from stream: {str(e)}", level="warning")
            return {}, True